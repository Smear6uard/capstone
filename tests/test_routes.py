from __future__ import annotations

from typing import Any
from io import BytesIO

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation

from app import main


@pytest.fixture(autouse=True)
def _clean_state() -> None:
    main._exam_configs.clear()
    main._exam_sessions.clear()
    main._completed_exams.clear()
    main._students.clear()
    main._student_ids_by_email.clear()
    main._material_chunks.clear()
    main._material_chunk_ids_by_config.clear()
    main._proctor_snapshots.clear()


@pytest.fixture
def client() -> TestClient:
    return TestClient(main.app)


def test_health_endpoints_return_ok(client: TestClient) -> None:
    for path in ("/", "/health"):
        response = client.get(path)

        assert response.status_code == 200
        assert response.json() == {"status": "ok"}


def test_configure_exam_persists_and_returns_config(client: TestClient) -> None:
    response = client.post(
        "/api/configure-exam",
        json={
            "domain": "Roman History",
            "topics": ["Fall of Rome", "Pax Romana"],
            "num_questions": 3,
            "difficulty": "medium",
            "special_instructions": "Emphasize primary sources.",
        },
    )

    assert response.status_code == 200
    payload = response.json()
    assert payload["domain"] == "Roman History"
    assert payload["topics"] == ["Fall of Rome", "Pax Romana"]
    assert payload["num_questions"] == 3
    assert payload["difficulty"] == "medium"
    assert payload["config_id"]

    listing = client.get("/api/exam-configs").json()
    assert len(listing) == 1
    assert listing[0]["config_id"] == payload["config_id"]


def test_start_exam_rejects_invalid_question_count(client: TestClient) -> None:
    response = client.post(
        "/api/start-exam",
        json={"domain": "Math", "num_questions": 99, "difficulty": "easy"},
    )

    assert response.status_code == 422


def test_generate_question_errors_when_session_missing(client: TestClient) -> None:
    response = client.post(
        "/api/generate-question",
        json={"session_id": "does-not-exist"},
    )

    assert response.status_code == 404


def test_full_exam_flow_grades_and_finishes(
    client: TestClient, monkeypatch: pytest.MonkeyPatch
) -> None:
    generate_call_count = {"count": 0}

    async def fake_call_together_json(**kwargs: Any) -> Any:
        model = kwargs["response_model"]
        if model is main.LLMGeneratedQuestion:
            generate_call_count["count"] += 1
            index = generate_call_count["count"]
            return main.LLMGeneratedQuestion(
                background_info=f"Context for question {index}.",
                question=f"Essay question #{index}?",
                grading_rubric=["Explains the topic", "Uses evidence"],
                topic=f"Subtopic {index}",
            )
        if model is main.GradeAnswerResponse:
            return main.GradeAnswerResponse(
                criterion_scores=[
                    main.CriterionScore(
                        criterion="Explains the topic",
                        score=82,
                        feedback="Add more specific examples.",
                    ),
                    main.CriterionScore(
                        criterion="Uses evidence",
                        score=88,
                        feedback="Good citations; cite one more source.",
                    ),
                ],
                overall_score=85,
                grading_explanation="Solid answer; add one concrete example.",
                question_index=0,
            )
        if model is main.LLMComposite:
            return main.LLMComposite(
                composite_score=86,
                composite_feedback="Strong work overall; deepen evidence.",
            )
        raise AssertionError(f"Unexpected response_model: {model}")

    monkeypatch.setattr(main, "call_together_json", fake_call_together_json)

    start_response = client.post(
        "/api/start-exam",
        json={
            "domain": "Roman History",
            "num_questions": 2,
            "difficulty": "medium",
            "topics": ["Fall of Rome", "Pax Romana"],
            "student_name": "Livia",
        },
    )
    assert start_response.status_code == 200
    session_id = start_response.json()["session_id"]

    for expected_index in range(2):
        q = client.post(
            "/api/generate-question", json={"session_id": session_id}
        )
        assert q.status_code == 200
        assert q.json()["question_index"] == expected_index
        assert q.json()["total_questions"] == 2

        g = client.post(
            "/api/grade-answer",
            json={
                "session_id": session_id,
                "student_answer": "My answer references political decline.",
                "time_spent_seconds": 300,
            },
        )
        assert g.status_code == 200
        assert g.json()["overall_score"] == 85

    finish = client.post("/api/finish-exam", json={"session_id": session_id})
    assert finish.status_code == 200
    body = finish.json()
    assert body["composite_score"] == 86
    assert body["num_questions"] == 2
    assert len(body["questions"]) == 2
    assert body["total_time_seconds"] == 600
    assert body["student_name"] == "Livia"

    results = client.get("/api/exam-results").json()
    assert len(results["exams"]) == 1
    assert results["exams"][0]["student_name"] == "Livia"


def test_student_login_and_history_tracks_multiple_completed_exams(
    client: TestClient, monkeypatch: pytest.MonkeyPatch
) -> None:
    grade_scores = [70, 90]
    generate_call_count = {"count": 0}

    async def fake_call_together_json(**kwargs: Any) -> Any:
        model = kwargs["response_model"]
        if model is main.LLMGeneratedQuestion:
            generate_call_count["count"] += 1
            return main.LLMGeneratedQuestion(
                background_info="Context.",
                question=f"Essay question {generate_call_count['count']}?",
                grading_rubric=["Explains the topic"],
                topic="Topic",
            )
        if model is main.GradeAnswerResponse:
            score = grade_scores.pop(0)
            return main.GradeAnswerResponse(
                criterion_scores=[
                    main.CriterionScore(
                        criterion="Explains the topic",
                        score=score,
                        feedback="Keep adding evidence.",
                    ),
                ],
                overall_score=score,
                grading_explanation="Good direction; add evidence.",
                question_index=0,
            )
        if model is main.LLMComposite:
            score = 70 if generate_call_count["count"] == 1 else 90
            return main.LLMComposite(
                composite_score=score,
                composite_feedback="Composite feedback.",
            )
        raise AssertionError(f"Unexpected response_model: {model}")

    monkeypatch.setattr(main, "call_together_json", fake_call_together_json)

    login = client.post(
        "/api/auth/login",
        json={"name": "Livia", "email": "LIVIA@example.edu"},
    )
    assert login.status_code == 200
    student_id = login.json()["student_id"]
    assert login.json()["email"] == "livia@example.edu"

    for domain in ("Roman History", "Greek History"):
        start = client.post(
            "/api/start-exam",
            json={
                "student_id": student_id,
                "domain": domain,
                "num_questions": 1,
                "difficulty": "medium",
            },
        )
        assert start.status_code == 200
        assert start.json()["student_name"] == "Livia"
        session_id = start.json()["session_id"]

        client.post("/api/generate-question", json={"session_id": session_id, "student_id": student_id})
        client.post(
            "/api/grade-answer",
            json={
                "session_id": session_id,
                "student_id": student_id,
                "student_answer": "A concise answer.",
                "time_spent_seconds": 60,
            },
        )
        finish = client.post(
            "/api/finish-exam",
            json={"session_id": session_id, "student_id": student_id},
        )
        assert finish.status_code == 200

    history = client.get(f"/api/student/history?student_id={student_id}")
    assert history.status_code == 200
    payload = history.json()
    assert payload["student_name"] == "Livia"
    assert payload["summary"]["total_exams"] == 2
    assert payload["summary"]["average_score"] == 80.0
    assert payload["summary"]["trend"] == "up"
    assert [session["domain"] for session in payload["sessions"]] == [
        "Greek History",
        "Roman History",
    ]
    assert payload["sessions"][0]["questions"][0]["effective_score"] == 90


def test_generate_question_rejects_when_exam_full(
    client: TestClient, monkeypatch: pytest.MonkeyPatch
) -> None:
    async def fake_call_together_json(**kwargs: Any) -> Any:
        model = kwargs["response_model"]
        if model is main.LLMGeneratedQuestion:
            return main.LLMGeneratedQuestion(
                background_info="Context.",
                question="Only essay question?",
                grading_rubric=["Clear explanation"],
                topic="Sole topic",
            )
        if model is main.GradeAnswerResponse:
            return main.GradeAnswerResponse(
                criterion_scores=[
                    main.CriterionScore(
                        criterion="Clear explanation",
                        score=70,
                        feedback="Organize the argument more clearly.",
                    ),
                ],
                overall_score=70,
                grading_explanation="Clear enough; tighten structure.",
                question_index=0,
            )
        raise AssertionError(f"Unexpected model {model}")

    monkeypatch.setattr(main, "call_together_json", fake_call_together_json)

    start = client.post(
        "/api/start-exam",
        json={"domain": "History", "num_questions": 1, "difficulty": "easy"},
    ).json()
    session_id = start["session_id"]

    client.post("/api/generate-question", json={"session_id": session_id})
    client.post(
        "/api/grade-answer",
        json={
            "session_id": session_id,
            "student_answer": "Something.",
            "time_spent_seconds": 60,
        },
    )

    second = client.post(
        "/api/generate-question", json={"session_id": session_id}
    )
    assert second.status_code == 400
    assert "already been generated" in second.json()["detail"]


def test_grade_answer_rejects_wrong_criterion_count(
    client: TestClient, monkeypatch: pytest.MonkeyPatch
) -> None:
    async def fake_call_together_json(**kwargs: Any) -> Any:
        model = kwargs["response_model"]
        if model is main.LLMGeneratedQuestion:
            return main.LLMGeneratedQuestion(
                background_info="Context.",
                question="Essay?",
                grading_rubric=["A", "B"],
                topic="Topic",
            )
        if model is main.GradeAnswerResponse:
            return main.GradeAnswerResponse(
                criterion_scores=[
                    main.CriterionScore(
                        criterion="A",
                        score=80,
                        feedback="Add evidence.",
                    ),
                ],
                overall_score=80,
                grading_explanation="Explanation.",
                question_index=0,
            )
        raise AssertionError

    monkeypatch.setattr(main, "call_together_json", fake_call_together_json)

    start = client.post(
        "/api/start-exam",
        json={"domain": "History", "num_questions": 1, "difficulty": "easy"},
    ).json()

    client.post("/api/generate-question", json={"session_id": start["session_id"]})

    graded = client.post(
        "/api/grade-answer",
        json={
            "session_id": start["session_id"],
            "student_answer": "Answer.",
            "time_spent_seconds": 60,
        },
    )

    assert graded.status_code == 502
    assert "wrong number of items" in graded.json()["detail"]


def test_finish_exam_rejects_no_graded_answers(client: TestClient) -> None:
    start = client.post(
        "/api/start-exam",
        json={"domain": "History", "num_questions": 1, "difficulty": "easy"},
    ).json()

    response = client.post(
        "/api/finish-exam", json={"session_id": start["session_id"]}
    )

    assert response.status_code == 400


def test_disputes_tutor_and_analytics_flow(
    client: TestClient, monkeypatch: pytest.MonkeyPatch
) -> None:
    generate_call_count = {"count": 0}
    grade_call_count = {"count": 0}
    tutor_call_count = {"count": 0}

    async def fake_call_together_json(**kwargs: Any) -> Any:
        model = kwargs["response_model"]
        if model is main.LLMGeneratedQuestion:
            generate_call_count["count"] += 1
            return main.LLMGeneratedQuestion(
                background_info="Context.",
                question=f"Essay question {generate_call_count['count']}?",
                grading_rubric=["Explains the cause", "Uses evidence"],
                topic="LLM-selected topic",
            )
        if model is main.GradeAnswerResponse:
            grade_call_count["count"] += 1
            score = 72 if grade_call_count["count"] == 1 else 90
            return main.GradeAnswerResponse(
                criterion_scores=[
                    main.CriterionScore(
                        criterion="Explains the cause",
                        score=score,
                        feedback="Connect the cause to the outcome.",
                    ),
                    main.CriterionScore(
                        criterion="Uses evidence",
                        score=score,
                        feedback="Use a more specific example.",
                    ),
                ],
                overall_score=score,
                grading_explanation="Good start; add evidence.",
                question_index=grade_call_count["count"] - 1,
            )
        if model is main.LLMComposite:
            return main.LLMComposite(
                composite_score=80,
                composite_feedback="Shows developing command of the material.",
            )
        if model is main.LLMGradeDisputeDecision:
            return main.LLMGradeDisputeDecision(
                dispute_accepted=True,
                revised_score=85,
                reviewer_explanation="The original answer included relevant evidence.",
            )
        if model is main.LLMTutorTurn:
            tutor_call_count["count"] += 1
            return main.LLMTutorTurn(
                message=f"Tutor turn {tutor_call_count['count']}: try one simpler step."
            )
        raise AssertionError(f"Unexpected response_model: {model}")

    monkeypatch.setattr(main, "call_together_json", fake_call_together_json)

    config = client.post(
        "/api/configure-exam",
        json={
            "domain": "Roman History",
            "topics": ["Pax Romana", "Fall of Rome"],
            "num_questions": 2,
            "difficulty": "medium",
            "grade_level": "High School",
            "grading_personality": "Strict",
            "teacher_name": "  Professor Elliott  ",
            "special_instructions": "Keep questions essay-based.",
        },
    )
    assert config.status_code == 200
    config_body = config.json()
    assert config_body["teacher_name"] == "Professor Elliott"
    assert config_body["grade_level"] == "High School"

    start = client.post(
        "/api/start-exam",
        json={
            "domain": "Ignored Domain",
            "num_questions": 1,
            "difficulty": "easy",
            "student_name": "Livia",
            "config_id": config_body["config_id"],
        },
    )
    assert start.status_code == 200
    start_body = start.json()
    session_id = start_body["session_id"]
    assert start_body["domain"] == "Roman History"
    assert start_body["num_questions"] == 2
    assert start_body["grade_level"] == "High School"
    assert start_body["grading_personality"] == "Strict"

    first = client.post("/api/generate-question", json={"session_id": session_id})
    assert first.status_code == 200
    assert first.json()["topic"] == "Pax Romana"
    client.post(
        "/api/grade-answer",
        json={
            "session_id": session_id,
            "student_answer": "The peace stabilized trade and governance.",
            "time_spent_seconds": 120,
        },
    )

    second = client.post("/api/generate-question", json={"session_id": session_id})
    assert second.status_code == 200
    assert second.json()["topic"] == "Fall of Rome"
    client.post(
        "/api/grade-answer",
        json={
            "session_id": session_id,
            "student_answer": "Military pressure and political instability mattered.",
            "time_spent_seconds": 180,
        },
    )

    finish = client.post("/api/finish-exam", json={"session_id": session_id})
    assert finish.status_code == 200
    finished = finish.json()
    assert finished["teacher_name"] == "Professor Elliott"
    assert finished["composite_feedback"].startswith("Professor Elliott's feedback:")

    dispute = client.post(
        "/api/dispute-grade",
        json={
            "session_id": session_id,
            "question_index": 0,
            "dispute_argument": "My answer did mention trade stability as evidence.",
        },
    )
    assert dispute.status_code == 200
    assert dispute.json() == {
        "dispute_accepted": True,
        "original_score": 72,
        "revised_score": 85,
        "reviewer_explanation": "The original answer included relevant evidence.",
    }

    duplicate_dispute = client.post(
        "/api/dispute-grade",
        json={
            "session_id": session_id,
            "question_index": 0,
            "dispute_argument": "Trying again.",
        },
    )
    assert duplicate_dispute.status_code == 400

    tutor_start = client.post(
        "/api/tutor-session",
        json={"session_id": session_id, "question_index": 0},
    )
    assert tutor_start.status_code == 200
    assert tutor_start.json()["messages"][0]["role"] == "tutor"

    tutor_reply = client.post(
        "/api/tutor-session",
        json={
            "session_id": session_id,
            "question_index": 0,
            "message": "I think stability made tax collection easier.",
        },
    )
    assert tutor_reply.status_code == 200
    assert [m["role"] for m in tutor_reply.json()["messages"]] == [
        "tutor",
        "student",
        "tutor",
    ]

    analytics = client.get("/api/exam-analytics")
    assert analytics.status_code == 200
    analytics_body = analytics.json()
    assert analytics_body["completed_sessions"] == 1
    assert analytics_body["sessions"][0]["dispute_count"] == 1
    assert analytics_body["sessions"][0]["composite_score"] == 86
    assert analytics_body["most_disputed_questions"][0]["accepted_disputes"] == 1
    assert analytics_body["per_question_average_scores"][0]["average_score"] == 85.0


def test_advanced_material_proctoring_and_adaptive_flow(
    client: TestClient, monkeypatch: pytest.MonkeyPatch
) -> None:
    grade_scores = [90, 92, 55, 50]
    generated_prompts: list[str] = []

    async def fake_call_together_json(**kwargs: Any) -> Any:
        model = kwargs["response_model"]
        user_content = kwargs["messages"][-1]["content"]
        if model is main.LLMGeneratedQuestion:
            generated_prompts.append(user_content)
            index = len(generated_prompts)
            return main.LLMGeneratedQuestion(
                background_info=f"Lecture-derived context {index}.",
                question=f"Explain lecture concept {index}?",
                grading_rubric=["Uses lecture evidence"],
                topic=f"Lecture Topic {index}",
            )
        if model is main.GradeAnswerResponse:
            score = grade_scores.pop(0)
            return main.GradeAnswerResponse(
                criterion_scores=[
                    main.CriterionScore(
                        criterion="Uses lecture evidence",
                        score=score,
                        feedback="Tie the answer to the uploaded slide.",
                    ),
                ],
                overall_score=score,
                grading_explanation="Good lecture-specific answer.",
                question_index=0,
            )
        if model is main.LLMComposite:
            return main.LLMComposite(
                composite_score=72,
                composite_feedback="Composite feedback.",
            )
        if model is main.LLMKnowledgeMap:
            return main.LLMKnowledgeMap(
                topics=[
                    main.KnowledgeMapEntry(
                        topic="Lecture Topic",
                        mastery_level="Proficient",
                        average_score=72,
                        summary="Understands the uploaded material unevenly.",
                    )
                ]
            )
        raise AssertionError(f"Unexpected response_model: {model}")

    async def fake_call_together_vision_json(**kwargs: Any) -> Any:
        return main.LLMProctorAnalysis(
            flags=["looking_away", "phone_visible"],
            confidence=0.82,
            description="Student appears to look down with a phone visible.",
        )

    monkeypatch.setattr(main, "call_together_json", fake_call_together_json)
    monkeypatch.setattr(main, "call_together_vision_json", fake_call_together_vision_json)

    config = client.post(
        "/api/configure-exam",
        json={
            "domain": "Cell Biology",
            "topics": ["Mitochondria"],
            "num_questions": 4,
            "difficulty": "medium",
        },
    ).json()

    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = "Mitochondria and ATP"
    textbox = slide.shapes.add_textbox(0, 0, 5_000_000, 1_000_000)
    textbox.text = "The lecture explains oxidative phosphorylation and ATP synthase."
    buf = BytesIO()
    deck.save(buf)
    buf.seek(0)

    upload = client.post(
        "/api/upload-material",
        data={"config_id": config["config_id"]},
        files={
            "file": (
                "lecture.pptx",
                buf.getvalue(),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )
    assert upload.status_code == 200
    chunks = upload.json()["chunks"]
    assert chunks[0]["source_label"] == "Slide 1"
    assert "oxidative phosphorylation" in chunks[0]["text"]

    edit = client.patch(
        "/api/material-chunk",
        json={
            "chunk_id": chunks[0]["chunk_id"],
            "text": chunks[0]["text"] + " The example compares ATP synthase to a turbine.",
        },
    )
    assert edit.status_code == 200

    login = client.post(
        "/api/auth/login",
        json={"name": "Mara Voss", "email": "mara@example.edu"},
    ).json()
    start = client.post(
        "/api/start-exam",
        json={
            "student_id": login["student_id"],
            "domain": "Ignored",
            "num_questions": 1,
            "difficulty": "easy",
            "config_id": config["config_id"],
        },
    ).json()
    session_id = start["session_id"]

    proctor = client.post(
        "/api/proctor/analyze",
        json={
            "session_id": session_id,
            "student_id": login["student_id"],
            "image_data_url": "data:image/jpeg;base64," + ("a" * 64),
        },
    )
    assert proctor.status_code == 200
    assert proctor.json()["confidence"] == 0.82

    seen_difficulties: list[str] = []
    for _ in range(4):
        question = client.post(
            "/api/generate-question",
            json={"session_id": session_id, "student_id": login["student_id"]},
        )
        assert question.status_code == 200
        seen_difficulties.append(question.json()["difficulty"])
        assert question.json()["source_label"] == "Slide 1"
        grade = client.post(
            "/api/grade-answer",
            json={
                "session_id": session_id,
                "student_id": login["student_id"],
                "student_answer": "This answer cites ATP synthase from the lecture.",
                "time_spent_seconds": 45,
            },
        )
        assert grade.status_code == 200

    assert seen_difficulties == ["medium", "medium", "hard", "hard"]
    assert "Uploaded material excerpts" in generated_prompts[0]

    finish = client.post(
        "/api/finish-exam",
        json={"session_id": session_id, "student_id": login["student_id"]},
    )
    assert finish.status_code == 200
    assert finish.json()["knowledge_map"][0]["mastery_level"] == "Proficient"

    alerts = client.get("/api/proctor/alerts")
    assert alerts.status_code == 200
    assert alerts.json()["sessions"][0]["integrity_score"] == 75
    assert alerts.json()["sessions"][0]["snapshots"][0]["flags"] == [
        "looking_away",
        "phone_visible",
    ]
